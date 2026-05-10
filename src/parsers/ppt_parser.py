from dataclasses import dataclass, field
from pathlib import Path
from pptx import Presentation
from pptx.util import Pt
from loguru import logger


@dataclass
class SlideChunk:
    text: str
    source: str        # 파일 경로
    slide_index: int   # 0-based
    slide_title: str
    doc_type: str = "ppt"
    metadata: dict = field(default_factory=dict)


def _extract_text_from_shape(shape) -> str:
    """도형에서 텍스트 추출. 테이블 포함."""
    if shape.has_text_frame:
        return "\n".join(p.text for p in shape.text_frame.paragraphs if p.text.strip())
    if shape.shape_type == 19:  # TABLE
        rows = []
        for row in shape.table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            rows.append(" | ".join(c for c in cells if c))
        return "\n".join(rows)
    return ""


def _get_slide_title(slide) -> str:
    """슬라이드 제목 추출. 없으면 빈 문자열."""
    if slide.shapes.title and slide.shapes.title.has_text_frame:
        return slide.shapes.title.text_frame.text.strip()
    return ""


def parse_pptx(file_path: str | Path) -> list[SlideChunk]:
    """
    PPTX 파일 파싱 → 슬라이드별 SlideChunk 리스트 반환.
    슬라이드 본문 + 스피커 노트를 합쳐 하나의 청크로 만든다.
    """
    path = Path(file_path)
    if not path.exists():
        raise FileNotFoundError(f"파일 없음: {path}")

    prs = Presentation(str(path))
    chunks: list[SlideChunk] = []

    for idx, slide in enumerate(prs.slides):
        title = _get_slide_title(slide)

        # 본문 텍스트 (제목 shape 제외)
        body_parts = []
        for shape in slide.shapes:
            if shape == slide.shapes.title:
                continue
            text = _extract_text_from_shape(shape)
            if text.strip():
                body_parts.append(text.strip())
        body = "\n".join(body_parts)

        # 스피커 노트
        notes = ""
        if slide.has_notes_slide:
            notes_tf = slide.notes_slide.notes_text_frame
            notes = notes_tf.text.strip() if notes_tf else ""

        # 제목 + 본문 + 노트 조합
        parts = [p for p in [title, body, notes] if p]
        full_text = "\n".join(parts)

        if not full_text.strip():
            logger.debug(f"  슬라이드 {idx+1}: 텍스트 없음 — 건너뜀")
            continue

        chunks.append(SlideChunk(
            text=full_text,
            source=str(path),
            slide_index=idx,
            slide_title=title,
            doc_type="ppt",
            metadata={"file_name": path.name, "total_slides": len(prs.slides)},
        ))

    logger.info(f"PPT 파싱 완료: {path.name} — {len(chunks)}/{len(prs.slides)} 슬라이드")
    return chunks


def parse_pptx_dir(dir_path: str | Path) -> list[SlideChunk]:
    """디렉토리 내 모든 PPTX 파일 일괄 파싱."""
    dir_path = Path(dir_path)
    files = list(dir_path.glob("**/*.pptx")) + list(dir_path.glob("**/*.PPTX"))
    if not files:
        logger.warning(f"PPTX 파일 없음: {dir_path}")
        return []

    all_chunks: list[SlideChunk] = []
    for f in files:
        try:
            all_chunks.extend(parse_pptx(f))
        except Exception as e:
            logger.error(f"파싱 실패 {f.name}: {e}")

    logger.info(f"총 {len(all_chunks)}개 슬라이드 청크 추출 ({len(files)}개 파일)")
    return all_chunks
